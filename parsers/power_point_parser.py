import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Iterator

from langchain_core.document_loaders import BaseBlobParser
from langchain_core.documents.base import Blob, Document


_ODP_MIME  = "application/vnd.oasis.opendocument.presentation"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PPT_MIME  = "application/vnd.ms-powerpoint"

# content_type hint forwarded to partition() to bypass libmagic misdetection
_MIME_HINTS: dict[str, str] = {
    ".pptx": _PPTX_MIME,
    ".ppt":  _PPT_MIME,
}


def _libreoffice_bin() -> str | None:
    return shutil.which("libreoffice") or shutil.which("soffice")


def _libreoffice_convert(source_path: str, target_ext: str) -> str:
    binary = _libreoffice_bin()
    if not binary:
        raise RuntimeError(
            "LibreOffice/soffice non trovato nel PATH: "
            "necessario per convertire file ODF"
        )
    src = Path(source_path)
    outdir = Path(tempfile.mkdtemp(prefix="odf-convert-"))
    try:
        proc = subprocess.run(
            [binary, "--headless", "--convert-to", target_ext,
             "--outdir", str(outdir), str(src)],
            capture_output=True, text=True,
        )
        if proc.returncode != 0:
            raise RuntimeError(
                f"Conversione LibreOffice fallita ({proc.returncode}): "
                f"stdout={proc.stdout!r} stderr={proc.stderr!r}"
            )
        converted = outdir / f"{src.stem}.{target_ext}"
        if not converted.exists():
            files = ", ".join(p.name for p in outdir.iterdir())
            raise RuntimeError(
                f"File convertito non trovato: atteso {converted.name}, "
                f"presenti: {files}"
            )
        return str(converted)
    except Exception:
        shutil.rmtree(outdir, ignore_errors=True)
        raise


def _extract_pptx_text(pptx_path: str) -> list[tuple[int, str]]:
    """Extract (slide_number, text) pairs directly via python-pptx.

    Avoids any import of unstructured.partition.pptx which depends on
    `unstructured.utils.lazyproperty` removed in unstructured >= 0.17.
    """
    from pptx import Presentation  # python-pptx is a direct dep of unstructured[pptx]

    prs = Presentation(pptx_path)
    results: list[tuple[int, str]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    chunks.append(line)
        if chunks:
            results.append((slide_num, "\n".join(chunks)))
    return results


class PowerPointParser(BaseBlobParser):
    """Lightweight PowerPoint/ODP parser used in non-multimodal mode.

    Strategy
    --------
    - .pptx  →  extract text directly via python-pptx (no unstructured import)
    - .ppt   →  LibreOffice headless → .pptx → python-pptx
    - .odp   →  LibreOffice headless → .pptx → python-pptx

    Rationale: unstructured.partition.pptx (and its transitive imports from
    unstructured.chunking / unstructured.utils) broke in unstructured >= 0.17
    due to the removal of `lazyproperty`.  python-pptx is always available
    because it is a direct dependency of `unstructured[pptx]`.

    NOTE: Blob.as_temp_file() is not available in all langchain-core versions;
    temp files are managed manually via tempfile.NamedTemporaryFile.
    """

    def lazy_parse(self, blob: Blob) -> Iterator[Document]:
        source = blob.source or ""
        suffix = os.path.splitext(source)[1].lower()
        temp_path: str | None = None
        converted_path: str | None = None
        converted_dir: str | None = None

        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                temp_path = tmp.name
                tmp.write(blob.as_bytes())

            if suffix in (".odp", ".ppt"):
                # ODP/PPT → convert to PPTX first
                converted_path = _libreoffice_convert(temp_path, "pptx")
                converted_dir = str(Path(converted_path).parent)
                pptx_path = converted_path
            else:
                # .pptx (or unknown) → use as-is
                pptx_path = temp_path

            slides = _extract_pptx_text(pptx_path)
            for slide_num, text in slides:
                metadata = blob.metadata.copy() if blob.metadata else {}
                metadata.update({
                    "source": source,
                    "element_type": "NarrativeText",
                    "page_number": slide_num,
                })
                yield Document(page_content=text, metadata=metadata)

        finally:
            for path in (temp_path, converted_path):
                if path and os.path.exists(path):
                    try:
                        os.unlink(path)
                    except Exception:
                        pass
            if converted_dir and os.path.isdir(converted_dir):
                shutil.rmtree(converted_dir, ignore_errors=True)
